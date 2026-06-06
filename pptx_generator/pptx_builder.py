import copy
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Cm, Pt

try:
    from .style_config import (
        BLOCK_COLORS,
        BLOCK_LAYOUT,
        BRANDING,
        COLORS,
        CONTENT_RULES,
        FONTS,
    )
except ImportError:
    from style_config import (
        BLOCK_COLORS,
        BLOCK_LAYOUT,
        BRANDING,
        COLORS,
        CONTENT_RULES,
        FONTS,
    )


BASE_DIR = Path(__file__).resolve().parent
ROOT_DIR = BASE_DIR.parent
DEFAULT_OUTPUT_PATH = BASE_DIR / "training_deck.pptx"
DEFAULT_RENDER_READY_PATH = BASE_DIR / "slides_render_ready.json"

SLIDE_W = 13.333
SLIDE_H = 7.5
FOOTER_H = 0.65


def hex_to_rgb(hex_color: str) -> RGBColor:
    value = (hex_color or "#000000").strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def color_for_key(key_or_hex: str) -> str:
    return COLORS.get(key_or_hex, key_or_hex)


def block_colors(block: str) -> dict[str, str]:
    config = BLOCK_COLORS.get(block, BLOCK_COLORS["theory"])
    return {
        "bg_color": color_for_key(config["bg"]),
        "title_color": color_for_key(config["title"]),
        "accent_color": color_for_key(config["accent"]),
    }


def normalize_notes(notes: dict[str, Any] | None, slide: dict[str, Any]) -> dict[str, str]:
    notes = notes or {}
    title = slide.get("title", "this slide")
    return {
        "aim": str(notes.get("aim") or f"Explain the purpose of {title}."),
        "time": str(notes.get("time") or "+/- 5 min"),
        "instructions": str(notes.get("instructions") or "Walk through the slide and invite one short reflection from the group."),
        "reflective_question": str(notes.get("reflective_question") or "How does this connect to your daily work?"),
        "debrief": str(notes.get("debrief") or "Summarize the key point before moving on."),
    }


def make_slide(
    block: str,
    title: str,
    bullets: list[str],
    notes: dict[str, str],
    layout_key: str | None = None,
) -> dict[str, Any]:
    colors = block_colors(block)
    return {
        "slide_index": 0,
        "block": block,
        "title": title,
        "bullets": bullets,
        "speaker_notes": notes,
        "layout_key": layout_key or BLOCK_LAYOUT.get(block, "title_body"),
        "bg_color": colors["bg_color"],
        "title_color": colors["title_color"],
        "accent_color": colors["accent_color"],
        "use_key_takeaway": block in CONTENT_RULES["key_takeaway_required_for"],
        "split_required": len(bullets) > CONTENT_RULES["max_bullets_per_slide"],
    }


def infer_training_title(slides: list[dict[str, Any]]) -> str:
    for slide in slides:
        if slide.get("block") == "cover" and slide.get("title"):
            return str(slide["title"])
    for slide in slides:
        if slide.get("title"):
            title = str(slide["title"])
            return title.replace("Welcome:", "").strip() or title
    return "Maverx Training"


def add_required_structural_slides(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    blocks = {slide.get("block") for slide in slides}
    title = infer_training_title(slides)
    additions: list[dict[str, Any]] = []

    if "cover" not in blocks:
        additions.append(make_slide(
            "cover",
            title,
            [
                "A practical training built around real participant tasks",
                "Clear concepts, recognizable examples, and active practice",
            ],
            {
                "aim": "Open the session and create a clear promise for participants.",
                "time": "2 minutes",
                "instructions": "Welcome the group and briefly connect the title to their work context.",
                "reflective_question": "What would make this session valuable for you?",
                "debrief": "Set the expectation that the training will stay practical and hands-on.",
            },
            "cover",
        ))

    if "agenda" not in blocks:
        additions.append(make_slide(
            "agenda",
            "What We Will Cover",
            ["Kick-off", "Theory", "Example", "Exercise", "Wrap-up"],
            {
                "aim": "Give participants a simple map of the session.",
                "time": "3 minutes",
                "instructions": "Walk through the flow and point out where participation is expected.",
                "reflective_question": "Which part are you most curious about?",
                "debrief": "Confirm the structure and transition into the kick-off.",
            },
            "agenda",
        ))

    if "timetable" not in blocks:
        additions.append(make_slide(
            "timetable",
            "Session Timeline",
            [
                "00:00 - Kick-off",
                "00:15 - Theory",
                "01:10 - Example",
                "01:40 - Exercise",
                "02:40 - Wrap-up",
            ],
            {
                "aim": "Help the trainer and participants understand the pacing.",
                "time": "2 minutes",
                "instructions": "Use this as a pacing reference and adapt if discussion runs long.",
                "reflective_question": "Does this timing fit the group energy today?",
                "debrief": "Keep the timeline visible as a trainer reference.",
            },
            "timetable",
        ))

    return additions + slides


def normalize_slides_for_rendering(slides_json: dict[str, Any]) -> dict[str, Any]:
    slides = copy.deepcopy(slides_json.get("slides", []))
    slides = add_required_structural_slides(slides)

    order = {
        "cover": 0,
        "agenda": 1,
        "timetable": 2,
        "kickoff": 3,
        "theory": 4,
        "example": 5,
        "exercise": 6,
        "debrief": 7,
        "wrapup": 8,
        "section": 9,
        "break": 10,
    }
    slides = sorted(enumerate(slides), key=lambda item: (order.get(item[1].get("block"), 99), item[0]))

    normalized = []
    for idx, (_, slide) in enumerate(slides):
        block = slide.get("block", "theory")
        colors = block_colors(block)
        slide["slide_index"] = idx
        slide["layout_key"] = slide.get("layout_key") or BLOCK_LAYOUT.get(block, "title_body")
        slide["bg_color"] = colors["bg_color"]
        slide["title_color"] = colors["title_color"]
        slide["accent_color"] = colors["accent_color"]
        slide["bullets"] = [str(b).lstrip("• ").strip() for b in slide.get("bullets", [])]
        slide["speaker_notes"] = normalize_notes(slide.get("speaker_notes"), slide)
        slide["use_key_takeaway"] = bool(slide.get("use_key_takeaway")) or block in CONTENT_RULES["key_takeaway_required_for"]
        slide["split_required"] = len(slide["bullets"]) > CONTENT_RULES["max_bullets_per_slide"]
        normalized.append(slide)

    return {"slides": normalized}


def add_textbox(slide, text: str, x: float, y: float, w: float, h: float, size: int, color: str, bold: bool = False):
    shape = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = FONTS["fallback"]
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = hex_to_rgb(color)
    return shape


def add_bullets(slide, bullets: list[str], x: float, y: float, w: float, h: float, color: str, size: int = 17):
    shape = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for i, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.text = str(bullet)
        paragraph.level = 0
        paragraph.space_after = Pt(7)
        paragraph.font.name = FONTS["fallback"]
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = hex_to_rgb(color)

    return shape


def add_footer(slide, slide_json: dict[str, Any]):
    bg = slide_json["bg_color"].upper()
    text_color = "#F2F2F2" if bg in {"#1A0040", "#3F0576"} else "#0D006A"
    accent = slide_json["accent_color"]

    line = slide.shapes.add_shape(1, Cm(0), Cm(SLIDE_H - FOOTER_H), Cm(SLIDE_W), Cm(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(accent)
    line.line.color.rgb = hex_to_rgb(accent)

    logo_path = BASE_DIR / BRANDING["logo_path"]
    if logo_path.exists():
        try:
            slide.shapes.add_picture(str(logo_path), Cm(0.55), Cm(SLIDE_H - 0.48), height=Cm(0.28))
        except Exception:
            add_textbox(slide, "Maverx", 0.55, SLIDE_H - 0.52, 2.0, 0.3, 8, text_color, True)
    else:
        add_textbox(slide, "Maverx", 0.55, SLIDE_H - 0.52, 2.0, 0.3, 8, text_color, True)

    site = add_textbox(slide, BRANDING["website"], SLIDE_W - 2.8, SLIDE_H - 0.52, 2.25, 0.3, 8, text_color)
    site.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def set_background(slide, color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color)


def render_cover(slide, slide_json: dict[str, Any]):
    add_textbox(slide, slide_json["title"], 0.9, 1.45, 11.2, 1.6, 34, slide_json["title_color"], True)
    add_bullets(slide, slide_json.get("bullets", [])[:3], 1.0, 3.35, 9.5, 1.55, "#F2F2F2", 18)
    bar = slide.shapes.add_shape(1, Cm(0.9), Cm(5.55), Cm(3.2), Cm(0.16))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(slide_json["accent_color"])
    bar.line.color.rgb = hex_to_rgb(slide_json["accent_color"])


def render_agenda(slide, slide_json: dict[str, Any]):
    add_title(slide, slide_json)
    bullets = slide_json.get("bullets", [])
    left = bullets[: (len(bullets) + 1) // 2]
    right = bullets[(len(bullets) + 1) // 2 :]
    add_bullets(slide, left, 1.1, 2.05, 5.2, 4.25, slide_json["title_color"], 18)
    add_bullets(slide, right, 7.0, 2.05, 5.0, 4.25, slide_json["title_color"], 18)


def render_timetable(slide, slide_json: dict[str, Any]):
    add_title(slide, slide_json)
    bullets = slide_json.get("bullets", [])
    y = 1.8
    for item in bullets[:7]:
        parts = str(item).split(" - ", 1)
        time = parts[0]
        activity = parts[1] if len(parts) > 1 else str(item)
        add_textbox(slide, time, 1.1, y, 2.0, 0.48, 14, slide_json["accent_color"], True)
        add_textbox(slide, activity, 3.35, y, 8.4, 0.48, 15, slide_json["title_color"])
        y += 0.64


def render_exercise(slide, slide_json: dict[str, Any]):
    add_title(slide, slide_json)
    bullets = slide_json.get("bullets", [])[:6]
    y = 1.75
    for i, bullet in enumerate(bullets, start=1):
        box = slide.shapes.add_shape(1, Cm(1.05), Cm(y), Cm(0.58), Cm(0.58))
        box.fill.solid()
        box.fill.fore_color.rgb = hex_to_rgb(slide_json["accent_color"])
        box.line.color.rgb = hex_to_rgb(slide_json["accent_color"])
        num = add_textbox(slide, str(i), 1.05, y + 0.05, 0.58, 0.35, 12, "#F2F2F2", True)
        num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        add_textbox(slide, str(bullet), 1.9, y - 0.02, 10.2, 0.62, 15, slide_json["title_color"])
        y += 0.72


def add_title(slide, slide_json: dict[str, Any]):
    add_textbox(slide, slide_json.get("title", ""), 0.85, 0.55, 11.3, 0.85, FONTS["title"], slide_json["title_color"], True)
    accent = slide.shapes.add_shape(1, Cm(0.85), Cm(1.43), Cm(1.25), Cm(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = hex_to_rgb(slide_json["accent_color"])
    accent.line.color.rgb = hex_to_rgb(slide_json["accent_color"])


def add_key_takeaway(slide, slide_json: dict[str, Any]):
    if not slide_json.get("use_key_takeaway"):
        return
    bullets = slide_json.get("bullets", [])
    takeaway = next((b for b in bullets if str(b).upper().startswith("KEY TAKEAWAY")), None)
    if not takeaway and bullets:
        takeaway = f"{CONTENT_RULES['key_takeaway_prefix']} {bullets[-1]}"
    if not takeaway:
        return

    box = slide.shapes.add_shape(1, Cm(0.85), Cm(6.2), Cm(11.65), Cm(0.58))
    box.fill.solid()
    box.fill.fore_color.rgb = hex_to_rgb(slide_json["accent_color"])
    box.line.color.rgb = hex_to_rgb(slide_json["accent_color"])
    add_textbox(slide, str(takeaway), 1.05, 6.29, 11.2, 0.32, 11, "#F2F2F2", True)


def render_standard(slide, slide_json: dict[str, Any]):
    add_title(slide, slide_json)
    bullets = slide_json.get("bullets", [])
    if slide_json.get("layout_key") == "title_2col" or len(bullets) > 5:
        left = bullets[: (len(bullets) + 1) // 2]
        right = bullets[(len(bullets) + 1) // 2 :]
        add_bullets(slide, left, 1.0, 1.85, 5.25, 3.95, slide_json["title_color"], 15)
        add_bullets(slide, right, 6.95, 1.85, 5.25, 3.95, slide_json["title_color"], 15)
    else:
        add_bullets(slide, bullets, 1.05, 1.85, 10.9, 4.25, slide_json["title_color"], 17)
    add_key_takeaway(slide, slide_json)


def add_speaker_notes(slide, notes: dict[str, str]):
    frame = slide.notes_slide.notes_text_frame
    frame.text = (
        f"AIM: {notes['aim']}\n"
        f"TIME: {notes['time']}\n"
        f"INSTRUCTIONS: {notes['instructions']}\n"
        f"REFLECTIVE QUESTION: {notes['reflective_question']}\n"
        f"DEBRIEF: {notes['debrief']}"
    )


def render_slide(prs: Presentation, slide_json: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, slide_json["bg_color"])

    layout = slide_json.get("layout_key")
    block = slide_json.get("block")
    if layout == "cover" or block == "cover":
        render_cover(slide, slide_json)
    elif layout == "agenda" or block == "agenda":
        render_agenda(slide, slide_json)
    elif layout == "timetable" or block == "timetable":
        render_timetable(slide, slide_json)
    elif layout == "exercise" or block == "exercise":
        render_exercise(slide, slide_json)
    else:
        render_standard(slide, slide_json)

    add_footer(slide, slide_json)
    add_speaker_notes(slide, slide_json["speaker_notes"])


def build_pptx(
    slides_json: dict[str, Any],
    output_path: str | Path = DEFAULT_OUTPUT_PATH,
    render_ready_path: str | Path = DEFAULT_RENDER_READY_PATH,
) -> Path:
    render_ready = normalize_slides_for_rendering(slides_json)
    render_ready_path = Path(render_ready_path)
    output_path = Path(output_path)

    render_ready_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with render_ready_path.open("w", encoding="utf-8") as f:
        json.dump(render_ready, f, indent=2, ensure_ascii=False)

    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)

    for slide_json in render_ready["slides"]:
        render_slide(prs, slide_json)

    # Remove the initial empty slide only if python-pptx created one.
    if len(prs.slides) > len(render_ready["slides"]):
        xml_slides = prs.slides._sldIdLst
        rel_id = xml_slides[0].rId
        prs.part.drop_rel(rel_id)
        xml_slides.remove(xml_slides[0])

    prs.save(output_path)
    return output_path


def main() -> None:
    input_path = BASE_DIR / "slides_pre_styled.json"
    with input_path.open("r", encoding="utf-8") as f:
        slides_json = json.load(f)

    output_path = build_pptx(slides_json)
    print(f"PPTX written to {output_path}")


if __name__ == "__main__":
    main()

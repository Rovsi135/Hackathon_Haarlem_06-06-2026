# to discover how to structure the json file for the pptx library
# run via
# python .\testing_pptx_library.py > output.txt

from pptx import Presentation

prs = Presentation("Maverx - Presentation Style Guide for Hackaton.pptx")


for i, slide in enumerate(prs.slides):
    print(f"\nSlide {i}: (layout: {slide.slide_layout.name})")
    for shape in slide.shapes:
        print(f"  shape: '{shape.name}' | has_text={shape.has_text_frame}")
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                print(f"    text: '{para.text}'")
                for run in para.runs:
                    try:
                        color = run.font.color.rgb
                    except AttributeError:
                        color = "scheme/inherited"
                    print(f"      run font: name={run.font.name}, size={run.font.size}, bold={run.font.bold}, color={color}")
        if shape.shape_type == 13:
            print(f"    [IMAGE]")